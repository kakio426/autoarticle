
import pytest
from pptx import Presentation
import os

def test_ppt_creation():
    from engines.ppt_engine import PPTEngine
    
    # Init
    engine = PPTEngine("웜 & 플레이풀", "Test School")
    
    # Create simple slide
    article = {
        "title": "Test Title",
        "content": "This is content.",
        "grade": "3rd Grade",
        "date": "2025-01-01",
        "location": "School", 
        "images": []
    }
    
    output_path = "test_output.pptx"
    if os.path.exists(output_path):
        os.remove(output_path)
        
    engine.create_presentation([article], output_path)
    
    assert os.path.exists(output_path)
    
    # Load and check basic structure
    prs = Presentation(output_path)
    assert len(prs.slides) >= 2 # Title slide + 1 content slide
    
    # Check Title Slide
    title_slide = prs.slides[0]
    # In new custom engine, we use Rectangle shapes, not default Title placeholders
    found_title = False
    for shape in title_slide.shapes:
        if shape.has_text_frame and "Test School 뉴스레터" in shape.text_frame.text:
            found_title = True
            break
    assert found_title, "Title text not found in custom shape"
    
    # Check Content Slide (Layout validation)
    content_slide = prs.slides[1]
    
    found_article_title = False
    for shape in content_slide.shapes:
        if shape.has_text_frame and "Test Title" in shape.text_frame.text:
            found_article_title = True
            # Check theme color (Orange-ish) applied to header text?
            # Actually header text is White, Background is Orange.
            # Let's just check text exists for now.
            break
    assert found_article_title, "Article title not found in content slide"
    
    os.remove(output_path)
